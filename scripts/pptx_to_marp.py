#!/usr/bin/env python3
"""Convert PowerPoint decks in this repository to Marp-compatible Markdown.

The converter favours clean, editable Markdown over pixel-perfect layout. It uses
python-pptx for slide text/images and reads the pptx package directly for speaker
notes, which are not currently exposed by python-pptx.
"""
from __future__ import annotations

import argparse
import html
import posixpath
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
NOTES_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"


@dataclass
class SlideContent:
    title: Optional[str] = None
    blocks: List[str] = field(default_factory=list)
    images: List[str] = field(default_factory=list)
    todos: List[str] = field(default_factory=list)
    notes: Optional[str] = None


def clean_text(text: str) -> str:
    """Normalise whitespace while preserving intentional line breaks."""
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def markdown_escape(text: str) -> str:
    # Keep escaping light so generated content remains readable.
    return text.replace("\u000b", "\n").strip()


def iter_shapes(shapes) -> Iterable:
    """Yield shapes, recursively descending into groups."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_sort_key(shape) -> Tuple[int, int]:
    return (int(getattr(shape, "top", 0) or 0), int(getattr(shape, "left", 0) or 0))


def is_title_shape(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        return shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        return False


def is_body_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        return shape.placeholder_format.type in (
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.VERTICAL_BODY,
        )
    except Exception:
        return False


def paragraph_is_bullet(paragraph) -> bool:
    ppr = paragraph._p.pPr
    if ppr is None:
        return paragraph.level > 0
    # Any bullet property except buNone indicates a bullet/numbered list.
    for child in ppr:
        local = child.tag.rsplit("}", 1)[-1]
        if local.startswith("bu") and local != "buNone":
            return True
    return paragraph.level > 0


def text_frame_to_markdown(tf, force_bullets: bool = False) -> str:
    lines: List[str] = []
    nonempty_count = sum(1 for paragraph in tf.paragraphs if clean_text(paragraph.text))
    for paragraph in tf.paragraphs:
        text = clean_text(paragraph.text)
        if not text:
            continue
        text = markdown_escape(text)
        if paragraph_is_bullet(paragraph) or (force_bullets and nonempty_count > 1):
            indent = "  " * max(0, paragraph.level)
            lines.append(f"{indent}- {text}")
        else:
            # Preserve explicit line breaks in paragraph text.
            lines.extend(text.splitlines())
    return "\n".join(lines).strip()


def table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        rows.append([clean_text(cell.text).replace("\n", " ") for cell in row.cells])
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * width) + " |"
    body = ["| " + " | ".join(row) + " |" for row in rows[1:]]
    return "\n".join([header, sep, *body]).strip()


def image_extension(shape) -> str:
    ext = getattr(shape.image, "ext", None) or "png"
    return ext.lower().lstrip(".")


def extract_slide(slide, slide_no: int, image_dir: Path, image_rel_dir: Path) -> SlideContent:
    content = SlideContent()
    text_shapes = []
    image_count = 0

    for shape in sorted(iter_shapes(slide.shapes), key=shape_sort_key):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            content.todos.append("Grouped layout was flattened; review ordering/layout.")
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            ext = image_extension(shape)
            filename = f"slide-{slide_no:02d}-image-{image_count:02d}.{ext}"
            path = image_dir / filename
            path.write_bytes(shape.image.blob)
            content.images.append(f"![]({(image_rel_dir / filename).as_posix()})")
            continue

        if getattr(shape, "has_chart", False):
            content.todos.append("Chart detected; convert or recreate manually.")
            continue

        if getattr(shape, "has_table", False):
            md_table = table_to_markdown(shape.table)
            if md_table:
                content.blocks.append(md_table)
            continue

        if getattr(shape, "has_text_frame", False):
            force_bullets = is_body_placeholder(shape) and not is_title_shape(shape)
            md = text_frame_to_markdown(shape.text_frame, force_bullets=force_bullets)
            if not md:
                continue
            if is_title_shape(shape) and content.title is None:
                content.title = re.sub(r"\s+", " ", md.replace("\n", " ")).strip()
            else:
                text_shapes.append(md)
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            continue

        if shape.shape_type not in (MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.AUTO_SHAPE):
            content.todos.append(f"Unsupported shape type {shape.shape_type}; review slide manually.")

    if content.title is None and text_shapes:
        # Promote the first short non-list text block to title when no title placeholder exists.
        first = text_shapes[0]
        if "\n" not in first and not first.lstrip().startswith("-") and len(first) <= 100:
            content.title = first
            text_shapes = text_shapes[1:]

    content.blocks.extend(text_shapes)
    # De-duplicate TODOs while preserving order.
    seen = set()
    content.todos = [t for t in content.todos if not (t in seen or seen.add(t))]
    return content


def resolve_target(base: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(base), target))


def notes_for_deck(pptx_path: Path) -> dict[int, str]:
    """Return speaker notes keyed by 1-based slide number."""
    notes: dict[int, str] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        for slide_no in range(1, 10_000):
            rels_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
            if rels_path not in zf.namelist():
                if slide_no > 1:
                    break
                continue
            rels_root = ET.fromstring(zf.read(rels_path))
            notes_target = None
            for rel in rels_root.findall("rel:Relationship", NS):
                if rel.attrib.get("Type") == NOTES_REL:
                    notes_target = rel.attrib.get("Target")
                    break
            if not notes_target:
                continue
            notes_path = resolve_target("ppt/slides/slide%d.xml" % slide_no, notes_target)
            if notes_path not in zf.namelist():
                continue
            root = ET.fromstring(zf.read(notes_path))
            texts: List[str] = []
            for shape in root.findall(".//p:sp", NS):
                ph = shape.find(".//p:ph", NS)
                if ph is not None and ph.attrib.get("type") in {"sldImg", "hdr", "ftr", "dt", "sldNum"}:
                    continue
                runs = [t.text or "" for t in shape.findall(".//a:t", NS)]
                text = clean_text(" ".join(runs))
                if text:
                    texts.append(text)
            if texts:
                notes[slide_no] = "\n\n".join(texts)
    return notes


def render_slide(content: SlideContent, slide_no: int) -> str:
    parts: List[str] = []
    title = content.title or f"Slide {slide_no}"
    parts.append(f"# {title}")
    for todo in content.todos:
        parts.append(f"<!-- TODO: {html.escape(todo)} -->")
    parts.extend(content.blocks)
    parts.extend(content.images)
    if content.notes:
        parts.append(f"<!--\n{content.notes}\n-->")
    return "\n\n".join(part for part in parts if part).strip()


def convert_deck(pptx_path: Path, slides_dir: Path) -> Path:
    deck_name = pptx_path.stem
    md_path = slides_dir / f"{deck_name}.md"
    image_dir = slides_dir / "images" / deck_name
    image_rel_dir = Path("images") / deck_name
    image_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    notes = notes_for_deck(pptx_path)
    rendered = ["---\nmarp: true\ntheme: default\npaginate: true\n---"]

    for idx, slide in enumerate(prs.slides, start=1):
        content = extract_slide(slide, idx, image_dir, image_rel_dir)
        content.notes = notes.get(idx)
        rendered.append(render_slide(content, idx))

    md_path.write_text("\n\n---\n\n".join(rendered) + "\n", encoding="utf-8")
    return md_path


def find_decks(root: Path) -> List[Path]:
    return sorted(
        path for path in root.rglob("*.pptx")
        if not path.name.startswith("~$") and ".git" not in path.parts and "slides" not in path.parts
    )


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert repository .pptx files to Marp Markdown.")
    parser.add_argument("decks", nargs="*", type=Path, help="Specific .pptx files to convert (default: all in repo).")
    parser.add_argument("--root", type=Path, default=Path.cwd(), help="Repository root to search from.")
    parser.add_argument("--slides-dir", type=Path, default=None, help="Output directory for Markdown and images.")
    args = parser.parse_args()

    root = args.root.resolve()
    slides_dir = (args.slides_dir or root / "slides").resolve()
    slides_dir.mkdir(parents=True, exist_ok=True)

    decks = [deck.resolve() for deck in args.decks] if args.decks else find_decks(root)
    if not decks:
        print("No .pptx files found.")
        return

    for deck in decks:
        md_path = convert_deck(deck, slides_dir)
        print(f"Converted {deck} -> {md_path}")


if __name__ == "__main__":
    main()
